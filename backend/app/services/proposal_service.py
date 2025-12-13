import os
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session
from app.models import RFP

OUTPUT_DIR = "/app/data/generated_proposals"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def generate_proposal_ppt(rfp_id: int, db: Session):
    rfp = db.query(RFP).filter(RFP.id == rfp_id).first()
    if not rfp: return {"error": "RFP not found"}
    
    data = rfp.extracted_data
    
    if not data:
        return {"error": "RFP data is empty. Please run the Technical Agent first."}
        
    if "line_items" not in data:
        return {"error": "Technical data missing (line_items). Please re-run Technical Analysis."}
        
    if "commercial" not in data:
        return {"error": "Pricing data missing. Please run Pricing Agent."}
    

    prs = Presentation()

    # --- SLIDE 1: TITLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Proposal for {rfp.client_name}"
    slide.placeholders[1].text = f"Ref: {rfp.title}\nBidWin AI - Multi-SKU Response"

    # --- SLIDE 2: SCOPE OF SUPPLY (TABLE) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
    slide.shapes.title.text = "Technical Scope & Matching"
    
    items = data["line_items"]
    rows = len(items) + 1
    cols = 4
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    
    headers = ["Requested Item", "Asian Paints Solution", "Score", "Reasoning"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h

    # Data
    for idx, item in enumerate(items):
        row = idx + 1
        match = item.get("match", {})
        scores = match.get("scores", {})
        
        table.cell(row, 0).text = item["requirement"].get("item_name", "N/A")
        table.cell(row, 1).text = match.get("product_name", "N/A")
        table.cell(row, 2).text = f"{scores.get('ensemble', 0)}%"
        table.cell(row, 3).text = match.get("reason", "")[:50] + "..." # Truncate

    # --- SLIDE 3: COMMERCIALS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Commercial Quote"
    
    comm = data["commercial"]
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"Total Project Value: ₹ {comm['grand_total_inr']}"
    
    p = tf.add_paragraph()
    p.text = "Includes Base Price, Logistics, Margin & GST."
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Detailed breakdown per item:"
    
    for line in comm["lines"]:
        p = tf.add_paragraph()
        p.text = f"- {line['item_name']} (Qty: {line['qty']}): ₹ {line['line_total']}"
        p.level = 1

    # Save
    filename = f"proposal_{rfp_id}.pptx"
    file_path = os.path.join(OUTPUT_DIR, filename)
    prs.save(file_path)
    
    rfp.status = "Ready to Submit"
    db.commit()

    return {
        "status": "success",
        "file_url": file_path,
        "download_url": f"/api/agents/main/download/{filename}"
    }